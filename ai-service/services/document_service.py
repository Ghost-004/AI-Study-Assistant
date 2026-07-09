from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import re
from io import BytesIO



def clean_text(text):
    # Collapse excessive blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Collapse multiple spaces
    text = re.sub(r"[ \t]+", " ", text)

    return text.strip()

def extract_text(contents, filename):
    extension = filename.lower().split(".")[-1]
    extractor = EXTRACTORS.get(extension)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {extension}")

    return extractor(contents)

def extract_txt(contents):
    return contents.decode("utf-8")

def extract_pptx(contents):
    ppt = Presentation(BytesIO(contents))
    text = ""

    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    text += shape.text + "\n"
    
    return text

def extract_pdf(contents):
    pdf = PdfReader(BytesIO(contents))
    text = ""

    for page in pdf.pages:
        extracted = page.extract_text()
        if extracted:
            text += extracted + "\n"
            
    return text

def extract_docx(contents):
    doc = Document(BytesIO(contents))
    text = ""

    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"

    return text

EXTRACTORS = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "pptx": extract_pptx,
    "txt": extract_txt,
}